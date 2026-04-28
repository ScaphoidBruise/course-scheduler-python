"""Generate slides.pptx for the COSC 3320 final-project demo.

Usage (run from project root):
    python presentation\\build_deck.py

The script:
  1. Builds a 10-slide 16:9 PowerPoint using python-pptx with the same color
     palette as scheduler/static/style.css for visual continuity with the app.
  2. Renders the architecture and database diagrams natively as PowerPoint
     shapes (no external image dependency). The .mmd source files in
     diagrams/ are kept as reference — paste them into mermaid.live if you
     ever want a generic mermaid render.
  3. Drops in screenshots from presentation/screenshots/ when present, or
     placeholders otherwise. See presentation/SCREENSHOTS.md.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


HERE = Path(__file__).resolve().parent
DIAGRAMS_DIR = HERE / "diagrams"
SCREENSHOTS_DIR = HERE / "screenshots"
ICONS_DIR = HERE / "icons"
OUT_PATH = HERE / "slides.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 10

# Color palette (matches scheduler/static/style.css)
NAVY_DEEP = RGBColor(0x0A, 0x16, 0x28)
NAVY_MID = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0xE8, 0x56, 0x2A)
ACCENT_HOVER = RGBColor(0xCF, 0x45, 0x20)
TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_BODY = RGBColor(0x2C, 0x3E, 0x52)
MUTED = RGBColor(0x88, 0x99, 0xAA)
SURFACE = RGBColor(0xF0, 0xF2, 0xF5)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x38, 0xA1, 0x69)
RED = RGBColor(0xE5, 0x3E, 0x3E)
GRID_LINE = RGBColor(0xE2, 0xE8, 0xF0)

FONT_SANS = "Segoe UI"
FONT_MONO = "Consolas"


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------

def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_solid_bg(slide, color: RGBColor) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = line_width
    return shape


def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font=FONT_SANS,
    size=18,
    bold=False,
    italic=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    v_anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(
    slide,
    paragraphs,
    left,
    top,
    width,
    height,
    *,
    font=FONT_SANS,
    size=16,
    color=TEXT_BODY,
    bullet="•",
    line_spacing=1.15,
    space_after_pt=8,
):
    """Each entry is dict(text=..., bold=..., color=..., size=..., bullet=..., indent=int)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    for i, item in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(item.get("space_after", space_after_pt))
        p.level = item.get("indent", 0)
        text = item["text"]
        b = item.get("bullet", bullet)
        if b:
            text = f"{b}  {text}"
        run = p.add_run()
        run.text = text
        run.font.name = item.get("font", font)
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.color.rgb = item.get("color", color)
    return tb


def add_accent_bar(slide, *, left=Inches(0.6), top=Inches(1.05), width=Inches(1.6), height=Inches(0.07)):
    add_rect(slide, left, top, width, height, ACCENT)


def add_title_header(slide, title: str, *, eyebrow: str | None = None):
    if eyebrow:
        add_textbox(
            slide,
            eyebrow.upper(),
            Inches(0.6),
            Inches(0.32),
            Inches(12),
            Inches(0.32),
            size=11,
            bold=True,
            color=ACCENT,
        )
        add_textbox(
            slide,
            title,
            Inches(0.6),
            Inches(0.55),
            Inches(12),
            Inches(0.7),
            size=30,
            bold=True,
            color=NAVY_DEEP,
        )
        add_accent_bar(slide, top=Inches(1.2))
    else:
        add_textbox(
            slide,
            title,
            Inches(0.6),
            Inches(0.4),
            Inches(12),
            Inches(0.7),
            size=30,
            bold=True,
            color=NAVY_DEEP,
        )
        add_accent_bar(slide, top=Inches(1.05))


def add_slide_number(slide, n: int):
    add_textbox(
        slide,
        f"{n} / {TOTAL_SLIDES}",
        Inches(12.5),
        Inches(7.05),
        Inches(0.7),
        Inches(0.3),
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    add_textbox(
        slide,
        "UTPB Course Scraper & Scheduler",
        Inches(0.6),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        size=10,
        color=MUTED,
    )


def add_placeholder(slide, label, left, top, width, height, *, fill=SURFACE, border=GRID_LINE):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT_SANS
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = MUTED
    return box


def add_image_or_placeholder(
    slide,
    image_path: Path | None,
    label: str,
    left,
    top,
    width,
    height,
):
    if image_path is not None and image_path.exists() and image_path.stat().st_size > 0:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    else:
        add_placeholder(slide, label, left, top, width, height)


def add_icon(slide, slug: str | None, left, top, size_w, size_h=None):
    """Drop a logo PNG from icons/ into the slide. No-op if slug is None or
    the file is missing/empty."""
    if not slug:
        return None
    path = ICONS_DIR / f"{slug}.png"
    if not path.exists() or path.stat().st_size == 0:
        return None
    if size_h is None:
        size_h = size_w
    return slide.shapes.add_picture(str(path), left, top, width=size_w, height=size_h)


def set_speaker_notes(slide, notes: str) -> None:
    nf = slide.notes_slide.notes_text_frame
    nf.text = notes


def add_card(slide, left, top, width, height, *, fill=CARD, border=GRID_LINE):
    return add_rect(slide, left, top, width, height, fill, line=border, line_width=Pt(0.75))


def add_pill(slide, label, left, top, width, height, *, fill=ACCENT, color=CARD):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT_SANS
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = color
    return box


# ---------------------------------------------------------------------------
# Native diagram primitives
# ---------------------------------------------------------------------------

def draw_arrow_right(slide, left, top, width, height, *, fill=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def draw_arrow_down(slide, left, top, width, height, *, fill=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def draw_line(slide, x1, y1, x2, y2, *, color=NAVY_MID, width_pt=1.25, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    if dashed:
        # python-pptx exposes dash via XML
        from lxml import etree
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return conn


def draw_lane(slide, left, top, width, height, header, items, *, header_fill=NAVY_DEEP,
              body_fill=CARD, item_color=TEXT_BODY, footer=None, footer_color=ACCENT,
              logo_strip=None, logo_strip_height=Inches(0.7)):
    """Swim-lane-style box: header + optional logo strip + items + optional footer.

    `items` entries can be plain strings, or (icon_slug, text) tuples to put a
    small icon to the left of the row text. `logo_strip` is a list of icon
    slugs rendered horizontally below the header bar.
    """
    add_card(slide, left, top, width, height, fill=body_fill, border=GRID_LINE)
    add_rect(slide, left, top, width, Inches(0.45), header_fill, line=header_fill)
    add_textbox(
        slide, header, left, top, width, Inches(0.45),
        size=12, bold=True, color=TEXT_LIGHT,
        align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
    )

    cursor_top = top + Inches(0.45)
    if logo_strip:
        n = len(logo_strip)
        max_icon = Emu(int(min(logo_strip_height - Inches(0.1), (width - Inches(0.4)) / max(n, 1))))
        gap = Inches(0.1)
        total_w = Emu(int(n * max_icon + (n - 1) * gap))
        start_x = left + Emu(int((width - total_w) / 2))
        icon_y = cursor_top + Emu(int((logo_strip_height - max_icon) / 2))
        for i, slug in enumerate(logo_strip):
            add_icon(slide, slug,
                     start_x + Emu(int(i * (max_icon + gap))),
                     icon_y, max_icon)
        cursor_top = cursor_top + logo_strip_height
    else:
        cursor_top = cursor_top + Inches(0.1)

    footer_h = Inches(0.32) if footer else Inches(0.0)
    available = top + height - cursor_top - footer_h
    if items:
        per_item = Emu(int(min(Inches(0.5), available / len(items))))
    else:
        per_item = Inches(0.4)

    for i, item in enumerate(items):
        row_top = cursor_top + Emu(int(i * per_item))
        if isinstance(item, tuple):
            icon_slug, text = item
            icon_size = Emu(int(min(per_item - Inches(0.08), Inches(0.22))))
            icon_left = left + Inches(0.08)
            icon_top = row_top + Emu(int((per_item - icon_size) / 2))
            add_icon(slide, icon_slug, icon_left, icon_top, icon_size)
            text_left = icon_left + icon_size + Inches(0.06)
            text_w = left + width - text_left - Inches(0.05)
            add_textbox(
                slide, text,
                text_left, row_top, text_w, per_item,
                size=9, color=item_color, font=FONT_MONO,
                align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
            )
        else:
            add_textbox(
                slide, item,
                left + Inches(0.05), row_top,
                width - Inches(0.1), per_item,
                size=10, color=item_color, font=FONT_MONO,
                align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
            )

    if footer:
        add_textbox(
            slide, footer,
            left, top + height - Inches(0.32),
            width, Inches(0.32),
            size=10, italic=True, color=footer_color,
            align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
        )


def draw_table_box(slide, left, top, width, height, name, columns,
                   *, header_fill=NAVY_DEEP, fk=False):
    """Draw an ER-style table box: name bar + 'col  TYPE' rows."""
    border_color = ACCENT if fk else GRID_LINE
    border_w = Pt(1.25) if fk else Pt(0.75)
    add_card(slide, left, top, width, height, fill=CARD, border=border_color)
    # Override line width
    add_rect(slide, left, top, width, Inches(0.32), header_fill, line=header_fill)
    add_textbox(
        slide, name, left, top, width, Inches(0.32),
        size=11, bold=True, color=TEXT_LIGHT, font=FONT_MONO,
        align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
    )
    rows_top = top + Inches(0.36)
    available = height - Inches(0.4)
    per_row = Emu(int(available / max(len(columns), 1)))
    for i, col in enumerate(columns):
        # col can be (name, type) or just a string
        if isinstance(col, tuple):
            label, kind = col
            row_text = f"{label}"
            kind_text = kind
        else:
            row_text = col
            kind_text = ""
        add_textbox(
            slide, row_text,
            left + Inches(0.12),
            rows_top + Emu(i) * per_row,
            width - Inches(1.0),
            per_row,
            size=9, color=TEXT_BODY, font=FONT_MONO,
            v_anchor=MSO_ANCHOR.MIDDLE,
        )
        if kind_text:
            add_textbox(
                slide, kind_text,
                left + width - Inches(1.0) - Inches(0.05),
                rows_top + Emu(i) * per_row,
                Inches(0.95),
                per_row,
                size=9, color=MUTED, font=FONT_MONO, align=PP_ALIGN.RIGHT,
                v_anchor=MSO_ANCHOR.MIDDLE,
            )


# ---------------------------------------------------------------------------
# Architecture diagram (slide 4)
# ---------------------------------------------------------------------------

def draw_architecture_diagram(slide, left, top, width, height):
    """Draw a 5-lane architecture diagram with arrows and an optional Anthropic box."""
    n_lanes = 5
    lane_w = Inches(1.5)
    arrow_gap = Emu(int((width - n_lanes * lane_w) / (n_lanes - 1)))
    lane_h = height - Inches(1.4)  # leave room for the optional Anthropic block

    lanes = [
        ("UTPB Sources", [
            ("site_smartcatalog", "SmartCatalog"),
            ("site_utpb", "Registrar"),
            ("site_utpb", "Academic Cal."),
            ("site_utpb", "Falcon Maps"),
        ], None, NAVY_DEEP, None),
        ("scrapers/", [
            "catalog.py",
            "sections.py",
            "session_dates.py",
            "infer_terms.py",
            "program_reqs.py",
        ], "stdlib only", NAVY_DEEP, ["python"]),
        ("data/courses.db", [
            "SQLite",
            "single file",
            "8 tables",
            "catalog + app data",
        ], "the hub", NAVY_MID, ["sqlite"]),
        ("scheduler/", [
            "app.py",
            "db.py",
            "conflict.py",
            "transcript_pdf.py",
        ], "Flask · JSON API", NAVY_DEEP, ["flask"]),
        ("Browser", [
            "pages/*.html",
            "static/*.js",
            "style.css",
            "vanilla JS only",
        ], None, NAVY_DEEP, ["html5", "css", "javascript"]),
    ]

    lane_lefts = []
    for i, (header, items, footer, header_fill, logos) in enumerate(lanes):
        lane_left = left + Emu(i) * (lane_w + arrow_gap)
        lane_lefts.append(lane_left)
        draw_lane(
            slide, lane_left, top, lane_w, lane_h,
            header, items,
            header_fill=header_fill, footer=footer, footer_color=ACCENT,
            logo_strip=logos,
        )

    # Arrows between lanes
    arrow_w = Inches(0.32)
    arrow_h = Inches(0.32)
    arrow_top = top + Emu(int(lane_h / 2)) - Emu(int(arrow_h / 2))
    for i in range(n_lanes - 1):
        gap_left = lane_lefts[i] + lane_w + Emu(int((arrow_gap - arrow_w) / 2))
        draw_arrow_right(slide, gap_left, arrow_top, arrow_w, arrow_h, fill=ACCENT)

    # Optional Anthropic API box, positioned below the Flask lane
    flask_left = lane_lefts[3]
    box_top = top + lane_h + Inches(0.4)
    box_w = lane_w + Emu(int(arrow_gap / 2))
    box_h = Inches(0.7)
    add_card(slide, flask_left, box_top, box_w, box_h,
             fill=SURFACE, border=MUTED)
    icon_size = Inches(0.45)
    icon_left = flask_left + Inches(0.1)
    icon_top = box_top + Emu(int((box_h - icon_size) / 2))
    has_icon = add_icon(slide, "anthropic", icon_left, icon_top, icon_size) is not None
    text_left = icon_left + icon_size + Inches(0.08) if has_icon else flask_left
    text_w = flask_left + box_w - text_left - Inches(0.05)
    add_textbox(
        slide, "Anthropic API",
        text_left, box_top + Inches(0.05), text_w, Inches(0.32),
        size=11, bold=True, color=NAVY_DEEP,
        align=PP_ALIGN.LEFT if has_icon else PP_ALIGN.CENTER,
        v_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, "optional · falls back to rules",
        text_left, box_top + Inches(0.35), text_w, Inches(0.3),
        size=9, italic=True, color=MUTED,
        align=PP_ALIGN.LEFT if has_icon else PP_ALIGN.CENTER,
        v_anchor=MSO_ANCHOR.MIDDLE,
    )

    # Dashed connector from Flask lane bottom to Anthropic box top
    flask_center_x = flask_left + Emu(int(lane_w / 2))
    draw_line(
        slide,
        flask_center_x, top + lane_h,
        flask_center_x, box_top,
        color=MUTED, width_pt=1.0, dashed=True,
    )


# ---------------------------------------------------------------------------
# Database diagram (slide 6)
# ---------------------------------------------------------------------------

def draw_database_diagram(slide, left, top, width, height):
    """Draw an 8-table ER diagram in a 3x3 grid (one cell intentionally blank)."""
    cols = 3
    rows = 3
    h_gap = Inches(0.22)
    v_gap = Inches(0.22)
    cell_w = Emu(int((width - (cols - 1) * h_gap) / cols))
    cell_h = Emu(int((height - (rows - 1) * v_gap) / rows))

    def cell_origin(col, row):
        return (
            left + Emu(col) * (cell_w + h_gap),
            top + Emu(row) * (cell_h + v_gap),
        )

    # Layout: (col, row) -> (table_name, columns, header_fill)
    tables = {
        (0, 0): ("users", [
            ("id", "PK"),
            ("username", "TEXT"),
            ("password_hash", "TEXT"),
            ("created_at", "TS"),
        ], NAVY_DEEP),
        (0, 1): ("user_profiles", [
            ("user_id", "PK/FK"),
            ("major", "TEXT"),
            ("minor", "TEXT"),
            ("transcript_json", "JSON"),
        ], NAVY_DEEP),
        (0, 2): ("course_wishlist", [
            ("id", "PK"),
            ("user_id", "FK"),
            ("course_id", "FK"),
        ], NAVY_DEEP),
        (1, 0): ("schedule_scenarios", [
            ("id", "PK"),
            ("user_id", "FK"),
            ("term", "TEXT"),
            ("name", "TEXT"),
            ("is_active", "BOOL"),
        ], NAVY_DEEP),
        (1, 1): ("user_schedules", [
            ("id", "PK"),
            ("scenario_id", "FK"),
            ("user_id", "FK"),
            ("section_id", "FK"),
        ], NAVY_DEEP),
        # (1, 2) intentionally blank
        (2, 0): ("courses", [
            ("id", "PK"),
            ("code", "TEXT"),
            ("name", "TEXT"),
            ("prereqs", "TEXT"),
            ("term_inferred", "TEXT"),
        ], NAVY_MID),
        (2, 1): ("sections", [
            ("id", "PK"),
            ("term", "TEXT"),
            ("section_code", "TEXT"),
            ("days/times", "TEXT"),
            ("session", "TEXT"),
        ], NAVY_MID),
        (2, 2): ("session_calendar", [
            ("term", "TEXT"),
            ("session", "TEXT"),
            ("start_date", "DATE"),
            ("end_date", "DATE"),
        ], NAVY_MID),
    }

    cell_box = {}  # (col, row) -> (left, top, right, bottom, cx, cy)
    for (col, row), (name, columns, fill) in tables.items():
        cl, ct = cell_origin(col, row)
        draw_table_box(slide, cl, ct, cell_w, cell_h, name, columns, header_fill=fill)
        cell_box[(col, row)] = (
            cl, ct, cl + cell_w, ct + cell_h,
            cl + Emu(int(cell_w / 2)),
            ct + Emu(int(cell_h / 2)),
        )

    def edge_point(box, side):
        cl, ct, cr, cb, cx, cy = box
        if side == "right":
            return (cr, cy)
        if side == "left":
            return (cl, cy)
        if side == "top":
            return (cx, ct)
        if side == "bottom":
            return (cx, cb)
        raise ValueError(side)

    # Relationship lines (from, to, side_from, side_to)
    rels = [
        ((0, 0), (0, 1), "bottom", "top"),  # users -> user_profiles
        ((0, 0), (1, 0), "right", "left"),  # users -> schedule_scenarios
        ((1, 0), (1, 1), "bottom", "top"),  # schedule_scenarios -> user_schedules
        ((1, 1), (2, 1), "right", "left"),  # user_schedules -> sections
        ((2, 0), (2, 1), "bottom", "top"),  # courses -> sections
        ((2, 1), (2, 2), "bottom", "top"),  # sections -> session_calendar
        ((0, 0), (0, 2), "left", "left"),   # users -> course_wishlist (run outside on the left)
    ]

    for src_pos, dst_pos, src_side, dst_side in rels:
        sx, sy = edge_point(cell_box[src_pos], src_side)
        dx, dy = edge_point(cell_box[dst_pos], dst_side)
        if src_side == "left" and dst_side == "left":
            # L-shape on the outside (users -> course_wishlist)
            offset = Inches(0.12)
            draw_line(slide, sx, sy, sx - offset, sy, color=NAVY_MID, width_pt=1.25)
            draw_line(slide, sx - offset, sy, dx - offset, dy, color=NAVY_MID, width_pt=1.25)
            draw_line(slide, dx - offset, dy, dx, dy, color=NAVY_MID, width_pt=1.25)
        else:
            draw_line(slide, sx, sy, dx, dy, color=NAVY_MID, width_pt=1.25)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_1_title(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, NAVY_DEEP)

    # Left text column
    add_rect(slide, Inches(0.6), Inches(2.6), Inches(0.08), Inches(2.4), ACCENT)

    add_textbox(
        slide,
        "COSC 3320 · FINAL PROJECT",
        Inches(0.85),
        Inches(2.6),
        Inches(7),
        Inches(0.4),
        size=12,
        bold=True,
        color=ACCENT,
    )
    add_textbox(
        slide,
        "UTPB Course Scraper",
        Inches(0.85),
        Inches(3.05),
        Inches(8),
        Inches(0.9),
        size=44,
        bold=True,
        color=TEXT_LIGHT,
    )
    add_textbox(
        slide,
        "& Scheduler",
        Inches(0.85),
        Inches(3.85),
        Inches(8),
        Inches(0.9),
        size=44,
        bold=True,
        color=ACCENT,
    )
    add_textbox(
        slide,
        "A local Python + Flask app that scrapes UTPB course data\nand turns it into a working term planner.",
        Inches(0.85),
        Inches(4.85),
        Inches(8),
        Inches(1.0),
        size=16,
        color=TEXT_LIGHT,
    )
    add_textbox(
        slide,
        "Your Name  ·  Spring 2026",
        Inches(0.85),
        Inches(6.4),
        Inches(8),
        Inches(0.4),
        size=14,
        color=MUTED,
    )

    # Right hero screenshot
    add_image_or_placeholder(
        slide,
        SCREENSHOTS_DIR / "hero_schedule.png",
        "Hero screenshot:\nSchedule page weekly grid",
        Inches(8.4),
        Inches(1.5),
        Inches(4.4),
        Inches(4.5),
    )

    set_speaker_notes(
        slide,
        "Hi, I'm <name>. For my COSC 3320 final project I built a local Python + Flask "
        "app that scrapes UTPB's catalog, registrar, and academic calendar and turns "
        "them into a working term planner. I'll show the architecture, the database, "
        "and then a short live demo. Total time: about 10 minutes.",
    )


def slide_2_motivation(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "Why this exists", eyebrow="Problem")

    bullets = [
        {"text": "UTPB's catalog, registrar sections, and academic calendar live on three different sites.", "bold": False},
        {"text": "Students manually juggle prerequisites, time conflicts, half-term sessions, and degree maps every term.", "bold": False},
        {"text": "Goal: one local app that pulls all that data and helps you build and validate a term plan.", "bold": True, "color": NAVY_DEEP},
    ]
    add_paragraphs(slide, bullets, Inches(0.6), Inches(1.6), Inches(7.6), Inches(4),
                   size=18, color=TEXT_BODY, space_after_pt=18)

    # Right "before / after" diagram with three sources collapsing into one app
    diagram_left = Inches(8.6)
    diagram_top = Inches(1.6)

    add_textbox(slide, "BEFORE", diagram_left, diagram_top, Inches(4.2), Inches(0.3),
                size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, diagram_left, diagram_top + Inches(0.35), Inches(4.2), Inches(0.6))
    add_textbox(slide, "SmartCatalog (courses + prereqs)", diagram_left, diagram_top + Inches(0.42),
                Inches(4.2), Inches(0.4), size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_card(slide, diagram_left, diagram_top + Inches(1.05), Inches(4.2), Inches(0.6))
    add_textbox(slide, "Registrar (live sections + times)", diagram_left, diagram_top + Inches(1.12),
                Inches(4.2), Inches(0.4), size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_card(slide, diagram_left, diagram_top + Inches(1.75), Inches(4.2), Inches(0.6))
    add_textbox(slide, "Academic calendar (session dates)", diagram_left, diagram_top + Inches(1.82),
                Inches(4.2), Inches(0.4), size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER)

    # Arrow down
    arrow_top = diagram_top + Inches(2.5)
    add_textbox(slide, "▼", diagram_left, arrow_top, Inches(4.2), Inches(0.4),
                size=18, color=ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(slide, "AFTER", diagram_left, arrow_top + Inches(0.45), Inches(4.2), Inches(0.3),
                size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, diagram_left, arrow_top + Inches(0.8), Inches(4.2), Inches(0.95),
             fill=NAVY_DEEP, border=NAVY_DEEP)
    add_textbox(slide, "One local app · one SQLite file", diagram_left,
                arrow_top + Inches(0.95), Inches(4.2), Inches(0.4),
                size=15, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "schedule · catalog · progress · planner", diagram_left,
                arrow_top + Inches(1.3), Inches(4.2), Inches(0.4),
                size=12, color=ACCENT, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 2)
    set_speaker_notes(
        slide,
        "Three different UTPB sites, one student trying to plan a semester. The motivation "
        "was simple: pull all of that data into one place and add the checks a human has "
        "to run by hand — conflicts, prereqs, half-term overlaps, and degree progress.",
    )


def slide_3_features(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "What it does", eyebrow="Feature tour")

    features = [
        ("Schedule builder",
         "Search/filter sections, drop them on a weekly grid. Conflict + prereq warnings.",
         "feature_schedule.png"),
        ("Catalog + wishlist",
         "Browse all courses with subject/level/term filters. Save courses to your wishlist.",
         "feature_catalog.png"),
        ("Transcript import",
         "Upload a UTPB PDF. Parsed in memory with pypdf — never written to disk.",
         "feature_profile.png"),
        ("Degree progress",
         "Completed / in-progress / remaining cards with a manual override flow.",
         "feature_progress.png"),
        ("Planner + AI advisor",
         "Multi-term roll-up, graduation forecast, .ics export, optional Anthropic advice.",
         "feature_planner.png"),
    ]

    # 5 cards in a single row across the slide
    card_top = Inches(1.6)
    card_h = Inches(5.0)
    margin = Inches(0.6)
    gutter = Inches(0.18)
    total_w = SLIDE_W - 2 * margin
    card_w = Emu(int((total_w - 4 * gutter) / 5))

    for i, (title, body, img_name) in enumerate(features):
        left = margin + Emu(i) * (card_w + gutter)
        add_card(slide, left, card_top, card_w, card_h)

        thumb_top = card_top + Inches(0.2)
        thumb_h = Inches(1.6)
        add_image_or_placeholder(
            slide,
            SCREENSHOTS_DIR / img_name,
            f"{title}\nscreenshot",
            left + Inches(0.18),
            thumb_top,
            card_w - Inches(0.36),
            thumb_h,
        )

        add_rect(slide, left + Inches(0.18), thumb_top + thumb_h + Inches(0.15),
                 Inches(0.6), Inches(0.04), ACCENT)

        add_textbox(
            slide,
            title,
            left + Inches(0.18),
            thumb_top + thumb_h + Inches(0.28),
            card_w - Inches(0.36),
            Inches(0.45),
            size=15,
            bold=True,
            color=NAVY_DEEP,
        )
        add_textbox(
            slide,
            body,
            left + Inches(0.18),
            thumb_top + thumb_h + Inches(0.78),
            card_w - Inches(0.36),
            Inches(2.0),
            size=11,
            color=TEXT_BODY,
        )

    add_slide_number(slide, 3)
    set_speaker_notes(
        slide,
        "Five features in one tour. The Schedule and Planner pages are the two I'll spend "
        "the most time on in the live demo. Catalog, Profile, and Progress all hang off "
        "the same SQLite database with a JSON API.",
    )


def slide_4_architecture(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "System architecture", eyebrow="How it fits together")

    draw_architecture_diagram(
        slide,
        Inches(0.6),
        Inches(1.55),
        Inches(8.6),
        Inches(5.2),
    )

    legend_left = Inches(9.4)
    legend_top = Inches(1.55)
    add_card(slide, legend_left, legend_top, Inches(3.4), Inches(5.2))
    add_textbox(slide, "Three layers, one DB", legend_left + Inches(0.2), legend_top + Inches(0.2),
                Inches(3.0), Inches(0.4), size=14, bold=True, color=NAVY_DEEP)

    legend = [
        {"text": "Scrapers feed SQLite — stdlib only, no requests / bs4.", "color": TEXT_BODY, "size": 12},
        {"text": "Flask reads SQLite, serves JSON + HTML.", "color": TEXT_BODY, "size": 12},
        {"text": "Vanilla JS pages call the JSON API.", "color": TEXT_BODY, "size": 12},
        {"text": "Anthropic call is optional and falls back to local rules.", "color": TEXT_BODY, "size": 12},
        {"text": "Data flows left → right; control stays on your machine.", "color": NAVY_DEEP, "size": 12, "bold": True},
    ]
    add_paragraphs(slide, legend,
                   legend_left + Inches(0.2), legend_top + Inches(0.7),
                   Inches(3.0), Inches(4.4), size=12, space_after_pt=10)

    add_slide_number(slide, 4)
    set_speaker_notes(
        slide,
        "Five scrapers feed one SQLite file. Flask reads from that same file, which keeps "
        "everything in one place. The frontend is plain HTML and vanilla JS — no React, no "
        "Jinja templating. The AI advisor is the only network dependency at runtime, and "
        "the endpoint falls back to rule-based advice when no API key is configured.",
    )


def slide_5_stack(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "Tech stack", eyebrow="Libraries & why")

    col_top = Inches(1.55)
    col_h = Inches(5.2)
    col_w = Inches(6.0)
    left_col = Inches(0.6)
    right_col = Inches(6.95)

    add_card(slide, left_col, col_top, col_w, col_h)
    add_card(slide, right_col, col_top, col_w, col_h)

    # Left: backend / scraping
    add_textbox(slide, "Backend & scraping", left_col + Inches(0.3), col_top + Inches(0.25),
                col_w, Inches(0.4), size=15, bold=True, color=NAVY_DEEP)
    add_rect(slide, left_col + Inches(0.3), col_top + Inches(0.7), Inches(0.6), Inches(0.04), ACCENT)

    left_items = [
        {"text": "Flask 3.1 — page routes + ~50 JSON endpoints in app.py", "bold": False},
        {"text": "Werkzeug — password hashing, secure session cookies", "bold": False},
        {"text": "sqlite3 (stdlib) — single file at data/courses.db", "bold": False},
        {"text": "pypdf — transcript parsed in memory, never saved to disk", "bold": False},
        {"text": "urllib + html + re + json — scrapers run on stdlib only", "bold": True, "color": NAVY_DEEP},
        {"text": "Anthropic Messages API — optional advisor, .env key", "bold": False, "color": MUTED},
    ]
    add_paragraphs(slide, left_items,
                   left_col + Inches(0.3), col_top + Inches(0.95),
                   col_w - Inches(0.6), Inches(4.2), size=13, space_after_pt=10)

    # Right: frontend / tooling
    add_textbox(slide, "Frontend & tooling", right_col + Inches(0.3), col_top + Inches(0.25),
                col_w, Inches(0.4), size=15, bold=True, color=NAVY_DEEP)
    add_rect(slide, right_col + Inches(0.3), col_top + Inches(0.7), Inches(0.6), Inches(0.04), ACCENT)

    right_items = [
        {"text": "Vanilla HTML + CSS + JS — no React/Vue/Jinja", "bold": True, "color": NAVY_DEEP},
        {"text": "Bootstrap utility classes for layout polish only", "bold": False},
        {"text": "Custom CSS palette in scheduler/static/style.css", "bold": False},
        {"text": "pytest / unittest — see tests/ for the suite", "bold": False},
        {"text": "python -m venv + pip — two-line setup in README", "bold": False},
        {"text": "git for version control; .env.example for secrets layout", "bold": False, "color": MUTED},
    ]
    add_paragraphs(slide, right_items,
                   right_col + Inches(0.3), col_top + Inches(0.95),
                   col_w - Inches(0.6), Inches(4.2), size=13, space_after_pt=10)

    # Bottom callout
    callout_top = col_top + col_h - Inches(0.85)
    add_rect(slide, left_col + Inches(0.3), callout_top, col_w - Inches(0.6), Inches(0.6),
             NAVY_DEEP)
    add_textbox(slide,
                "Why stdlib scrapers? Portable, deterministic, zero install pain on a grader's machine.",
                left_col + Inches(0.4), callout_top + Inches(0.1),
                col_w - Inches(0.8), Inches(0.45),
                size=11, italic=True, color=TEXT_LIGHT)
    add_rect(slide, right_col + Inches(0.3), callout_top, col_w - Inches(0.6), Inches(0.6),
             NAVY_DEEP)
    add_textbox(slide,
                "Why no JS framework? Course is graded on Python; the frontend stays small enough to read.",
                right_col + Inches(0.4), callout_top + Inches(0.1),
                col_w - Inches(0.8), Inches(0.45),
                size=11, italic=True, color=TEXT_LIGHT)

    add_slide_number(slide, 5)
    set_speaker_notes(
        slide,
        "Two short lists. The most interesting choice on the left is that the scrapers use "
        "the standard library only — no requests, no BeautifulSoup. On the right, no JS "
        "framework: the frontend is intentionally small so a Python-class grader can read it.",
    )


def slide_6_database(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "Database schema", eyebrow="One file · data/courses.db")

    draw_database_diagram(
        slide,
        Inches(0.6),
        Inches(1.55),
        Inches(8.6),
        Inches(5.2),
    )

    legend_left = Inches(9.4)
    legend_top = Inches(1.55)
    add_card(slide, legend_left, legend_top, Inches(3.4), Inches(5.2))
    add_textbox(slide, "8 tables, 1 SQLite file",
                legend_left + Inches(0.2), legend_top + Inches(0.2),
                Inches(3.0), Inches(0.4), size=14, bold=True, color=NAVY_DEEP)

    legend = [
        {"text": "Catalog: courses + sections + session_calendar", "color": TEXT_BODY, "size": 12},
        {"text": "Auth: users (Werkzeug-hashed passwords)", "color": TEXT_BODY, "size": 12},
        {"text": "Profile: user_profiles holds parsed transcript JSON", "color": TEXT_BODY, "size": 12},
        {"text": "Schedules: schedule_scenarios + user_schedules per term", "color": TEXT_BODY, "size": 12},
        {"text": "Wishlist: course_wishlist", "color": TEXT_BODY, "size": 12},
        {"text": "Scrapers update only their tables — accounts and saved schedules survive a re-scrape.",
         "color": NAVY_DEEP, "size": 12, "bold": True},
    ]
    add_paragraphs(slide, legend,
                   legend_left + Inches(0.2), legend_top + Inches(0.7),
                   Inches(3.0), Inches(4.4), size=12, space_after_pt=10)

    add_slide_number(slide, 6)
    set_speaker_notes(
        slide,
        "All eight tables live in one SQLite file. Catalog data and app data co-exist on "
        "purpose — backups are a single file copy. The catalog scraper only touches the "
        "courses table, so re-running scrapers never wipes a user account or saved schedule.",
    )


def slide_7_demo(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, NAVY_DEEP)

    add_textbox(slide, "LIVE DEMO", Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
                size=12, bold=True, color=ACCENT)
    add_textbox(slide, "Switching to the running app", Inches(0.6), Inches(0.85),
                Inches(12), Inches(0.7), size=30, bold=True, color=TEXT_LIGHT)
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(1.6), Inches(0.07), ACCENT)

    steps = [
        ("1", "Schedule",   "Pick a term, search/filter, add a section, trigger a conflict."),
        ("2", "Profile",    "Upload a transcript PDF — parsed in memory only."),
        ("3", "Progress",   "Completed / in-progress / remaining + manual override."),
        ("4", "Planner",    "Multi-term roll-up, graduation forecast, AI advisor button."),
        ("5", "Export .ics", "Download the active schedule and open it in a real calendar."),
    ]

    step_top = Inches(2.0)
    for i, (num, name, desc) in enumerate(steps):
        row_top = step_top + Inches(0.95 * i)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.6), row_top, Inches(0.7), Inches(0.7))
        circle.shadow.inherit = False
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.name = FONT_SANS
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TEXT_LIGHT

        add_textbox(slide, name, Inches(1.5), row_top + Inches(0.0),
                    Inches(3.0), Inches(0.45),
                    size=20, bold=True, color=TEXT_LIGHT)
        add_textbox(slide, desc, Inches(1.5), row_top + Inches(0.42),
                    Inches(11), Inches(0.4),
                    size=14, color=MUTED)

    add_slide_number(slide, 7)
    set_speaker_notes(
        slide,
        "I'll be on the app for the next 4 minutes. The path is on the slide as a safety "
        "net — if anything goes wrong I jump to the next step. Watch for the conflict "
        "warning on Schedule and the graduation forecast on Planner; those are the two "
        "moments that show the data is actually wired up.",
    )


def slide_8_design(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "Python design highlights", eyebrow="The interesting bits")

    # Two cards side by side
    card_top = Inches(1.55)
    card_h = Inches(5.2)
    card_w = Inches(6.0)
    left_col = Inches(0.6)
    right_col = Inches(6.95)

    # Left: half-semester aware conflict detection
    add_card(slide, left_col, card_top, card_w, card_h)
    add_textbox(slide, "Half-semester aware conflict detection",
                left_col + Inches(0.3), card_top + Inches(0.25),
                card_w - Inches(0.6), Inches(0.45),
                size=15, bold=True, color=NAVY_DEEP)
    add_textbox(slide, "scheduler/conflict.py",
                left_col + Inches(0.3), card_top + Inches(0.7),
                card_w - Inches(0.6), Inches(0.3),
                size=10, italic=True, color=MUTED, font=FONT_MONO)

    code_box = add_card(slide, left_col + Inches(0.3), card_top + Inches(1.05),
                        card_w - Inches(0.6), Inches(2.55),
                        fill=NAVY_DEEP, border=NAVY_DEEP)
    code_lines = [
        {"text": "if is_half_semester(a['session']) \\", "color": ACCENT, "bullet": ""},
        {"text": "    and is_half_semester(b['session']):", "color": ACCENT, "bullet": ""},
        {"text": "    a_w1 = 'W1' in a['session']", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "    b_w1 = 'W1' in b['session']", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "    if a_w1 != b_w1:", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "        return False  # 8W1 and 8W2", "color": GREEN, "bullet": ""},
        {"text": "                       # never overlap", "color": GREEN, "bullet": ""},
        {"text": "return start_a < end_b and start_b < end_a", "color": TEXT_LIGHT, "bullet": ""},
    ]
    add_paragraphs(slide, code_lines,
                   left_col + Inches(0.45), card_top + Inches(1.2),
                   card_w - Inches(0.9), Inches(2.3),
                   font=FONT_MONO, size=11, color=TEXT_LIGHT,
                   line_spacing=1.2, space_after_pt=2)

    add_textbox(slide,
                "Two 8-week sessions in the same term don't actually run at the same time, "
                "so the app refuses to false-flag them.",
                left_col + Inches(0.3), card_top + Inches(3.8),
                card_w - Inches(0.6), Inches(1.2),
                size=12, color=TEXT_BODY)

    # Right: AI advisor with graceful fallback
    add_card(slide, right_col, card_top, card_w, card_h)
    add_textbox(slide, "AI advisor with graceful fallback",
                right_col + Inches(0.3), card_top + Inches(0.25),
                card_w - Inches(0.6), Inches(0.45),
                size=15, bold=True, color=NAVY_DEEP)
    add_textbox(slide, "scheduler/app.py · _call_ai_planner_advisor",
                right_col + Inches(0.3), card_top + Inches(0.7),
                card_w - Inches(0.6), Inches(0.3),
                size=10, italic=True, color=MUTED, font=FONT_MONO)

    add_card(slide, right_col + Inches(0.3), card_top + Inches(1.05),
             card_w - Inches(0.6), Inches(2.55),
             fill=NAVY_DEEP, border=NAVY_DEEP)
    code_lines2 = [
        {"text": "api_key = os.environ.get('ANTHROPIC_API_KEY')", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "if not api_key:", "color": ACCENT, "bullet": ""},
        {"text": "    return None  # caller uses local rules", "color": GREEN, "bullet": ""},
        {"text": "", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "payload = {", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "  'model': model,", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "  'messages': _messages_for(context),", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "  'max_tokens': 450,", "color": TEXT_LIGHT, "bullet": ""},
        {"text": "}", "color": TEXT_LIGHT, "bullet": ""},
    ]
    add_paragraphs(slide, code_lines2,
                   right_col + Inches(0.45), card_top + Inches(1.2),
                   card_w - Inches(0.9), Inches(2.3),
                   font=FONT_MONO, size=11, color=TEXT_LIGHT,
                   line_spacing=1.2, space_after_pt=2)

    add_textbox(slide,
                "No key, no problem: the endpoint switches to rule-based advice so the "
                "feature still works in a live demo or on a grader's machine.",
                right_col + Inches(0.3), card_top + Inches(3.8),
                card_w - Inches(0.6), Inches(1.2),
                size=12, color=TEXT_BODY)

    add_slide_number(slide, 8)
    set_speaker_notes(
        slide,
        "Two design choices I want to call out. First, conflict detection knows that the "
        "two 8-week sessions don't overlap, so I don't get bogus warnings when stacking a "
        "first-half and second-half class. Second, the AI advisor degrades gracefully — "
        "no key just means the user gets local rule-based advice. That keeps the demo "
        "deterministic.",
    )


def slide_9_testing(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "Testing & data safety", eyebrow="Quality")

    # Left: coverage list
    left_col = Inches(0.6)
    col_top = Inches(1.55)
    col_w = Inches(6.0)
    col_h = Inches(5.2)
    add_card(slide, left_col, col_top, col_w, col_h)
    add_textbox(slide, "What the suite covers",
                left_col + Inches(0.3), col_top + Inches(0.25),
                col_w, Inches(0.4),
                size=15, bold=True, color=NAVY_DEEP)
    add_rect(slide, left_col + Inches(0.3), col_top + Inches(0.7),
             Inches(0.6), Inches(0.04), ACCENT)

    items = [
        {"text": "tests/test_prereqs.py — prereq parser ('SUBJ ####' forms, AND/OR)"},
        {"text": "tests/test_planner_api.py — overview, target credits, term timeline"},
        {"text": "tests/test_scenarios.py — create / duplicate / rename / delete / activate"},
        {"text": "tests/test_progress_overview.py — degree-progress summary endpoint"},
        {"text": "tests/test_account_api.py — register / login / change password / delete"},
        {"text": "tests/test_program_requirements.py — scraped program review tables"},
    ]
    add_paragraphs(slide, items,
                   left_col + Inches(0.3), col_top + Inches(0.95),
                   col_w - Inches(0.6), Inches(3.5),
                   size=12, color=TEXT_BODY,
                   font=FONT_MONO, space_after_pt=8)

    # Run-command callout
    cmd_top = col_top + Inches(4.4)
    add_rect(slide, left_col + Inches(0.3), cmd_top, col_w - Inches(0.6), Inches(0.55),
             NAVY_DEEP)
    add_textbox(slide,
                "python -m pytest tests/ -v",
                left_col + Inches(0.45), cmd_top + Inches(0.1),
                col_w - Inches(0.9), Inches(0.4),
                size=14, font=FONT_MONO, bold=True, color=ACCENT)

    # Right: pytest screenshot + data-safety blurb
    right_col = Inches(6.95)
    add_card(slide, right_col, col_top, col_w, col_h)
    add_textbox(slide, "Last green run",
                right_col + Inches(0.3), col_top + Inches(0.25),
                col_w, Inches(0.4),
                size=15, bold=True, color=NAVY_DEEP)
    add_rect(slide, right_col + Inches(0.3), col_top + Inches(0.7),
             Inches(0.6), Inches(0.04), ACCENT)

    add_image_or_placeholder(
        slide,
        SCREENSHOTS_DIR / "pytest_output.png",
        "Screenshot:\n`python -m pytest tests/ -v`\n(green run)",
        right_col + Inches(0.3), col_top + Inches(0.95),
        col_w - Inches(0.6), Inches(2.6),
    )

    safety_top = col_top + Inches(3.7)
    add_textbox(slide, "Data safety",
                right_col + Inches(0.3), safety_top,
                col_w, Inches(0.4),
                size=13, bold=True, color=NAVY_DEEP)
    safety_items = [
        {"text": "--backup-db copies the whole SQLite file before a catalog refresh."},
        {"text": "Scrapers update only their own tables — accounts, schedules, profiles stay put."},
        {"text": "Transcript PDFs are parsed in memory and discarded; only the parsed JSON is stored."},
    ]
    add_paragraphs(slide, safety_items,
                   right_col + Inches(0.3), safety_top + Inches(0.4),
                   col_w - Inches(0.6), Inches(1.4),
                   size=11, color=TEXT_BODY, space_after_pt=4)

    add_slide_number(slide, 9)
    set_speaker_notes(
        slide,
        "Six test files cover the parts of the app most likely to break silently — prereq "
        "parsing, the planner overview, scenarios, degree progress, accounts, and "
        "program-requirement scraping. The right side calls out the data-safety guarantees: "
        "scrapers don't wipe the database, and transcripts never touch disk.",
    )


def slide_10_wrap(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    set_solid_bg(slide, SURFACE)
    add_title_header(slide, "Wins, limits, what's next", eyebrow="Wrap")

    # Three columns: wins / limits / next
    col_top = Inches(1.55)
    col_h = Inches(4.4)
    margin = Inches(0.6)
    gutter = Inches(0.25)
    col_w = Emu(int((SLIDE_W - 2 * margin - 2 * gutter) / 3))

    columns = [
        ("What worked", GREEN, [
            "One SQLite file = trivial backups, easy grading.",
            "No JS framework — frontend stayed readable.",
            "Stdlib scrapers survived a registrar redesign mid-semester.",
            "Conflict detection respects 8-week sessions.",
        ]),
        ("Limits", RED, [
            "Scrapers are HTML-shape sensitive; selectors break on redesigns.",
            "Prereq parser only handles 'SUBJ ####' patterns — free text is treated as unparseable.",
            "No mobile layout pass yet.",
            "AI advisor uses a single prompt template — no chat memory.",
        ]),
        ("Next", ACCENT, [
            "Switch scrapers to a small HTML-tolerant parser layer.",
            "Add iCal subscription URL, not just .ics export.",
            "Cache program requirements per program version.",
            "Mobile-first restyle of Schedule and Planner.",
        ]),
    ]

    for i, (heading, color, bullets) in enumerate(columns):
        left = margin + Emu(i) * (col_w + gutter)
        add_card(slide, left, col_top, col_w, col_h)
        add_rect(slide, left, col_top, col_w, Inches(0.08), color)
        add_textbox(slide, heading, left + Inches(0.25), col_top + Inches(0.25),
                    col_w, Inches(0.4),
                    size=16, bold=True, color=NAVY_DEEP)
        add_paragraphs(
            slide,
            [{"text": b} for b in bullets],
            left + Inches(0.25),
            col_top + Inches(0.75),
            col_w - Inches(0.5),
            col_h - Inches(0.9),
            size=12,
            color=TEXT_BODY,
            space_after_pt=8,
            line_spacing=1.2,
        )

    # Bottom strip: thanks + repo (kept above the slide-number footer at 7.05")
    strip_top = col_top + col_h + Inches(0.2)
    strip_h = Inches(0.7)
    add_rect(slide, Inches(0.6), strip_top, Inches(12.13), strip_h, NAVY_DEEP)
    add_textbox(slide, "Thanks — questions?",
                Inches(0.85), strip_top + Inches(0.12),
                Inches(7), Inches(0.5),
                size=20, bold=True, color=TEXT_LIGHT)
    add_textbox(slide, "github.com/<you>/utpb-scheduler",
                Inches(7.5), strip_top + Inches(0.22),
                Inches(5.2), Inches(0.4),
                size=14, color=ACCENT, font=FONT_MONO, align=PP_ALIGN.RIGHT)

    add_slide_number(slide, 10)
    set_speaker_notes(
        slide,
        "Three columns: what worked, what didn't, and what I'd do next. The recurring theme "
        "is that I traded scrape robustness for transparency — every scraper is short "
        "enough to read in a sitting. The biggest follow-up is a tolerant HTML parser layer "
        "so a registrar redesign doesn't take the app down. Happy to take questions.",
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    DIAGRAMS_DIR.mkdir(parents=True, exist_ok=True)
    SCREENSHOTS_DIR.mkdir(parents=True, exist_ok=True)
    print("Building slides …")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_1_title,
        slide_2_motivation,
        slide_3_features,
        slide_4_architecture,
        slide_5_stack,
        slide_6_database,
        slide_7_demo,
        slide_8_design,
        slide_9_testing,
        slide_10_wrap,
    ]
    for fn in builders:
        fn(prs)
        print(f"  + {fn.__name__}")

    prs.save(OUT_PATH)
    print(f"\nWrote {OUT_PATH.relative_to(HERE.parent)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
