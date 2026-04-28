"""Quick structural sanity check for slides.pptx."""
from pptx import Presentation
from pptx.util import Emu

p = Presentation('presentation/slides.pptx')
print(f'Slides: {len(p.slides)}  ({Emu(p.slide_width).inches}x{Emu(p.slide_height).inches} in)')

for i, s in enumerate(p.slides, 1):
    text_runs = []
    n_shapes = 0
    n_pics = 0
    n_text = 0
    for sh in s.shapes:
        n_shapes += 1
        if getattr(sh, 'shape_type', None) is not None and sh.shape_type == 13:
            n_pics += 1
        if sh.has_text_frame:
            n_text += 1
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.text.strip():
                        text_runs.append(r.text.strip())
    blob = ' | '.join(text_runs)[:200]
    notes = ''
    if s.has_notes_slide:
        nt = s.notes_slide.notes_text_frame.text
        notes = (nt[:80] + '…') if len(nt) > 80 else nt
    print(f'  Slide {i}: shapes={n_shapes} text={n_text} pics={n_pics}')
    print(f'    text:  {blob}')
    print(f'    notes: {notes}')
